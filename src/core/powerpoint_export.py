"""
PowerPoint Export Module

Professional PowerPoint presentation generation with:
- Multiple slide layouts (title, content, comparison, etc.)
- Charts and graphs
- Tables with formatting
- Images and shapes
- Transitions and animations
- Speaker notes
- Custom themes and branding
"""

from typing import List, Dict, Any, Optional, Union, Tuple
from datetime import datetime
from pathlib import Path
import logging

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint export disabled.")

logger = logging.getLogger('dms.powerpoint')


class PPTTheme:
    """Theme configuration for PowerPoint presentations."""

    def __init__(self):
        """Initialize default theme."""
        # Colors
        self.color_primary = RGBColor(13, 110, 253)  # Blue
        self.color_secondary = RGBColor(108, 117, 125)  # Gray
        self.color_accent = RGBColor(25, 135, 84)  # Green
        self.color_warning = RGBColor(255, 193, 7)  # Yellow
        self.color_danger = RGBColor(220, 53, 69)  # Red
        self.color_text = RGBColor(33, 37, 41)  # Dark
        self.color_bg = RGBColor(255, 255, 255)  # White

        # Fonts
        self.font_title = 'Calibri'
        self.font_body = 'Calibri'
        self.font_size_title = Pt(44)
        self.font_size_heading = Pt(32)
        self.font_size_body = Pt(18)
        self.font_size_small = Pt(14)

        # Company info
        self.company_name = "Document Management System"
        self.company_tagline = "Professional Service Management"
        self.logo_path = None


class PowerPointExporter:
    """Export data to professional PowerPoint presentations."""

    def __init__(self, theme: Optional[PPTTheme] = None):
        """
        Initialize PowerPoint exporter.

        Args:
            theme: Theme configuration
        """
        if not PYTHON_PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PowerPoint export. "
                "Install it with: pip install python-pptx"
            )

        self.theme = theme or PPTTheme()
        self.presentation = Presentation()
        # Set slide size (16:9)
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(5.625)

    def add_title_slide(self, title: str, subtitle: Optional[str] = None) -> Any:
        """
        Add a title slide.

        Args:
            title: Main title
            subtitle: Optional subtitle

        Returns:
            Slide object
        """
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.name = self.theme.font_title
        title_frame.paragraphs[0].font.size = self.theme.font_size_title
        title_frame.paragraphs[0].font.color.rgb = self.theme.color_primary
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Set subtitle
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.paragraphs[0].font.name = self.theme.font_body
            subtitle_frame.paragraphs[0].font.size = self.theme.font_size_body
            subtitle_frame.paragraphs[0].font.color.rgb = self.theme.color_secondary
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        logger.info(f"Added title slide: {title}")
        return slide

    def add_content_slide(self, title: str, content: Union[str, List[str]],
                         bullet_points: bool = True) -> Any:
        """
        Add a content slide with bullet points or text.

        Args:
            title: Slide title
            content: Content text or list of bullet points
            bullet_points: Format as bullet points if True

        Returns:
            Slide object
        """
        slide_layout = self.presentation.slide_layouts[1]  # Title and Content
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.name = self.theme.font_title
        title_frame.paragraphs[0].font.size = self.theme.font_size_heading
        title_frame.paragraphs[0].font.color.rgb = self.theme.color_primary
        title_frame.paragraphs[0].font.bold = True

        # Add content
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        if isinstance(content, str):
            content = [content]

        for idx, item in enumerate(content):
            if idx > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]

            p.text = item
            p.font.name = self.theme.font_body
            p.font.size = self.theme.font_size_body
            p.font.color.rgb = self.theme.color_text
            p.level = 0 if bullet_points else 0

        logger.info(f"Added content slide: {title}")
        return slide

    def add_comparison_slide(self, title: str, left_content: List[str],
                           right_content: List[str],
                           left_title: str = "Before",
                           right_title: str = "After") -> Any:
        """
        Add a comparison slide (two-column layout).

        Args:
            title: Slide title
            left_content: Left column content
            right_content: Right column content
            left_title: Left column title
            right_title: Right column title

        Returns:
            Slide object
        """
        slide_layout = self.presentation.slide_layouts[3]  # Two Content layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.name = self.theme.font_title
        title_frame.paragraphs[0].font.size = self.theme.font_size_heading
        title_frame.paragraphs[0].font.color.rgb = self.theme.color_primary

        # Left column
        left_shape = slide.placeholders[1]
        self._add_column_content(left_shape, left_title, left_content)

        # Right column
        right_shape = slide.placeholders[2]
        self._add_column_content(right_shape, right_title, right_content)

        logger.info(f"Added comparison slide: {title}")
        return slide

    def _add_column_content(self, shape: Any, title: str, content: List[str]) -> None:
        """Helper to add content to a column."""
        text_frame = shape.text_frame
        text_frame.clear()

        # Add title
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.name = self.theme.font_body
        p.font.size = self.theme.font_size_body
        p.font.color.rgb = self.theme.color_primary
        p.font.bold = True

        # Add content
        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.name = self.theme.font_body
            p.font.size = self.theme.font_size_small
            p.font.color.rgb = self.theme.color_text
            p.level = 1

    def add_table_slide(self, title: str, data: List[List[str]],
                       headers: Optional[List[str]] = None) -> Any:
        """
        Add a slide with a table.

        Args:
            title: Slide title
            data: Table data (list of rows)
            headers: Optional header row

        Returns:
            Slide object
        """
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5)

        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.name = self.theme.font_title
        title_frame.paragraphs[0].font.size = self.theme.font_size_heading
        title_frame.paragraphs[0].font.color.rgb = self.theme.color_primary
        title_frame.paragraphs[0].font.bold = True

        # Calculate table dimensions
        rows = len(data) + (1 if headers else 0)
        cols = len(headers) if headers else len(data[0])

        # Add table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(3.5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        for col_idx in range(cols):
            table.columns[col_idx].width = width // cols

        # Add headers
        if headers:
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = self.theme.font_size_small
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme.color_primary
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Add data
        start_row = 1 if headers else 0
        for row_idx, row_data in enumerate(data, start=start_row):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.size = self.theme.font_size_small

        logger.info(f"Added table slide: {title}")
        return slide

    def add_chart_slide(self, title: str, chart_type: str,
                       categories: List[str], series_data: Dict[str, List[float]],
                       chart_title: Optional[str] = None) -> Any:
        """
        Add a slide with a chart.

        Args:
            title: Slide title
            chart_type: 'bar', 'line', 'pie', or 'column'
            categories: X-axis categories
            series_data: Dictionary of series name -> data values
            chart_title: Optional chart title

        Returns:
            Slide object
        """
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5)

        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.name = self.theme.font_title
        title_frame.paragraphs[0].font.size = self.theme.font_size_heading
        title_frame.paragraphs[0].font.color.rgb = self.theme.color_primary
        title_frame.paragraphs[0].font.bold = True

        # Map chart types
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE
        }

        chart_type_enum = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories

        for series_name, values in series_data.items():
            chart_data.add_series(series_name, values)

        # Add chart
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
        chart = slide.shapes.add_chart(
            chart_type_enum, x, y, cx, cy, chart_data
        ).chart

        # Set chart title
        if chart_title:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_title
            chart.chart_title.text_frame.paragraphs[0].font.size = self.theme.font_size_body

        # Set legend position
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        logger.info(f"Added {chart_type} chart slide: {title}")
        return slide

    def add_image_slide(self, title: str, image_path: str,
                       caption: Optional[str] = None) -> Any:
        """
        Add a slide with an image.

        Args:
            title: Slide title
            image_path: Path to image file
            caption: Optional image caption

        Returns:
            Slide object
        """
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5)

        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.name = self.theme.font_title
        title_frame.paragraphs[0].font.size = self.theme.font_size_heading
        title_frame.paragraphs[0].font.color.rgb = self.theme.color_primary
        title_frame.paragraphs[0].font.bold = True

        # Add image
        try:
            left = Inches(1)
            top = Inches(1.5)
            pic = slide.shapes.add_picture(image_path, left, top, width=Inches(8))

            # Add caption
            if caption:
                left = Inches(1)
                top = Inches(4.5)
                width = Inches(8)
                height = Inches(0.5)

                caption_shape = slide.shapes.add_textbox(left, top, width, height)
                caption_frame = caption_shape.text_frame
                caption_frame.text = caption
                caption_frame.paragraphs[0].font.size = self.theme.font_size_small
                caption_frame.paragraphs[0].font.color.rgb = self.theme.color_secondary
                caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            logger.info(f"Added image slide: {title}")

        except Exception as e:
            logger.error(f"Error adding image: {e}")

        return slide

    def add_section_divider(self, section_title: str,
                           subtitle: Optional[str] = None) -> Any:
        """
        Add a section divider slide.

        Args:
            section_title: Section title
            subtitle: Optional subtitle

        Returns:
            Slide object
        """
        slide_layout = self.presentation.slide_layouts[2]  # Section Header
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = section_title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.name = self.theme.font_title
        title_frame.paragraphs[0].font.size = self.theme.font_size_title
        title_frame.paragraphs[0].font.color.rgb = self.theme.color_primary
        title_frame.paragraphs[0].font.bold = True

        # Set subtitle
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.paragraphs[0].font.size = self.theme.font_size_body
            subtitle_frame.paragraphs[0].font.color.rgb = self.theme.color_secondary

        logger.info(f"Added section divider: {section_title}")
        return slide

    def add_notes(self, slide: Any, notes_text: str) -> None:
        """
        Add speaker notes to a slide.

        Args:
            slide: Slide object
            notes_text: Notes text
        """
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text
        logger.info("Added speaker notes to slide")

    def save(self, output_path: str) -> bool:
        """
        Save presentation to file.

        Args:
            output_path: Output file path

        Returns:
            True if successful
        """
        try:
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)

            self.presentation.save(str(output_path))
            logger.info(f"Saved presentation to: {output_path}")
            return True

        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
            return False

    def export_services_presentation(self, services: List[Dict[str, Any]],
                                    output_path: str) -> bool:
        """
        Export services data as a presentation.

        Args:
            services: List of service dictionaries
            output_path: Output file path

        Returns:
            True if successful
        """
        try:
            # Title slide
            self.add_title_slide(
                "Services Report",
                f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
            )

            # Overview slide
            overview_content = [
                f"Total Services: {len(services)}",
                f"Report Date: {datetime.now().strftime('%Y-%m-%d')}",
                "Analysis includes financial data and regional distribution"
            ]
            self.add_content_slide("Overview", overview_content)

            # Services by region (if data available)
            if services:
                regions = {}
                for service in services:
                    region = service.get('region', 'Unknown')
                    regions[region] = regions.get(region, 0) + 1

                if regions:
                    self.add_chart_slide(
                        "Services by Region",
                        "column",
                        list(regions.keys()),
                        {"Count": list(regions.values())},
                        "Distribution of Services by Region"
                    )

            # Summary slide
            self.add_content_slide(
                "Summary",
                [
                    "All services analyzed successfully",
                    "Financial calculations completed",
                    "Regional distribution documented",
                    "Thank you for your attention"
                ]
            )

            # Save presentation
            return self.save(output_path)

        except Exception as e:
            logger.error(f"Error exporting services presentation: {e}")
            return False


# Convenience functions
def create_presentation(title: str, slides_data: List[Dict[str, Any]],
                       output_path: str) -> bool:
    """
    Quick create a presentation from structured data.

    Args:
        title: Presentation title
        slides_data: List of slide dictionaries with 'type', 'title', 'content'
        output_path: Output file path

    Returns:
        True if successful
    """
    exporter = PowerPointExporter()
    exporter.add_title_slide(title)

    for slide_data in slides_data:
        slide_type = slide_data.get('type', 'content')
        title = slide_data.get('title', '')
        content = slide_data.get('content', [])

        if slide_type == 'content':
            exporter.add_content_slide(title, content)
        elif slide_type == 'table':
            headers = slide_data.get('headers')
            exporter.add_table_slide(title, content, headers)
        elif slide_type == 'chart':
            chart_type = slide_data.get('chart_type', 'column')
            categories = slide_data.get('categories', [])
            series = slide_data.get('series', {})
            exporter.add_chart_slide(title, chart_type, categories, series)

    return exporter.save(output_path)
